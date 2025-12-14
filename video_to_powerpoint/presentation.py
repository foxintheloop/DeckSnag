"""PowerPoint presentation management."""

import io
import logging
from pathlib import Path
from typing import List, Optional
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger("video_to_powerpoint")


class PresentationManager:
    """Manages PowerPoint presentation creation and slide addition.

    Unlike the original implementation, this keeps the presentation in memory
    and only saves to disk when explicitly requested, improving performance.
    """

    # Standard slide dimensions (16:9 aspect ratio)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Image positioning
    IMAGE_LEFT = Inches(0.5)
    IMAGE_TOP = Inches(0.5)
    IMAGE_MAX_WIDTH = Inches(12.333)
    IMAGE_MAX_HEIGHT = Inches(6.5)

    def __init__(self) -> None:
        """Initialize the presentation manager."""
        self._presentation: Optional[Presentation] = None
        self._slides: List[Image.Image] = []
        self._output_path: Optional[Path] = None

    def create(self, output_path: Optional[Path] = None) -> None:
        """Create a new presentation.

        Args:
            output_path: Optional path to save the presentation.
        """
        self._presentation = Presentation()

        # Set slide dimensions to 16:9
        self._presentation.slide_width = self.SLIDE_WIDTH
        self._presentation.slide_height = self.SLIDE_HEIGHT

        self._slides = []
        self._output_path = Path(output_path) if output_path else None

        logger.info("Created new presentation")

    def add_slide(self, image: Image.Image) -> int:
        """Add a slide with an image to the presentation.

        Args:
            image: PIL Image to add as a slide.

        Returns:
            The slide number (1-indexed).

        Raises:
            RuntimeError: If no presentation has been created.
        """
        if self._presentation is None:
            raise RuntimeError("No presentation created. Call create() first.")

        # Store the image
        self._slides.append(image.copy())

        # Get blank slide layout (index 6 is typically blank)
        blank_layout = self._presentation.slide_layouts[6]

        # Add slide
        slide = self._presentation.slides.add_slide(blank_layout)

        # Convert image to bytes for pptx
        img_bytes = io.BytesIO()
        image.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        # Calculate image dimensions to fit slide while maintaining aspect ratio
        img_width, img_height = image.size
        aspect_ratio = img_width / img_height

        max_width = self.IMAGE_MAX_WIDTH
        max_height = self.IMAGE_MAX_HEIGHT

        # Determine dimensions
        if img_width / max_width.inches > img_height / max_height.inches:
            # Width is the limiting factor
            width = max_width
            height = Inches(max_width.inches / aspect_ratio)
        else:
            # Height is the limiting factor
            height = max_height
            width = Inches(max_height.inches * aspect_ratio)

        # Center the image on the slide
        left = Inches((self.SLIDE_WIDTH.inches - width.inches) / 2)
        top = Inches((self.SLIDE_HEIGHT.inches - height.inches) / 2)

        # Add image to slide
        slide.shapes.add_picture(img_bytes, left, top, width=width, height=height)

        slide_num = len(self._slides)
        logger.info(f"Added slide {slide_num}")

        return slide_num

    def save(self, path: Optional[Path] = None) -> Path:
        """Save the presentation to disk.

        Args:
            path: Path to save to. Uses the path from create() if not specified.

        Returns:
            Path where the presentation was saved.

        Raises:
            RuntimeError: If no presentation exists or no path specified.
        """
        if self._presentation is None:
            raise RuntimeError("No presentation to save. Call create() first.")

        save_path = Path(path) if path else self._output_path

        if save_path is None:
            raise RuntimeError("No output path specified")

        # Ensure the file has .pptx extension
        if save_path.suffix.lower() != ".pptx":
            save_path = save_path.with_suffix(".pptx")

        # Ensure parent directory exists
        save_path.parent.mkdir(parents=True, exist_ok=True)

        # Save the presentation
        self._presentation.save(str(save_path))
        logger.info(f"Saved presentation to: {save_path}")

        return save_path

    def get_slide_count(self) -> int:
        """Get the number of slides in the presentation.

        Returns:
            Number of slides.
        """
        return len(self._slides)

    def get_slides(self) -> List[Image.Image]:
        """Get copies of all slide images.

        Returns:
            List of PIL Images.
        """
        return [img.copy() for img in self._slides]

    def get_slide(self, index: int) -> Optional[Image.Image]:
        """Get a specific slide image by index.

        Args:
            index: 0-based slide index.

        Returns:
            PIL Image or None if index is invalid.
        """
        if 0 <= index < len(self._slides):
            return self._slides[index].copy()
        return None

    def get_last_slide(self) -> Optional[Image.Image]:
        """Get the most recently added slide image.

        Returns:
            PIL Image of last slide, or None if no slides exist.
        """
        if self._slides:
            return self._slides[-1].copy()
        return None

    @property
    def has_slides(self) -> bool:
        """Check if the presentation has any slides."""
        return len(self._slides) > 0

    def close(self) -> None:
        """Close the presentation and clear resources."""
        self._presentation = None
        self._slides = []
        self._output_path = None
        logger.debug("Presentation closed")
